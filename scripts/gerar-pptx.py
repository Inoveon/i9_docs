# -*- coding: utf-8 -*-
"""
Script para gerar apresentação PowerPoint da i9ON
Execute: pip install python-pptx && python gerar_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# Cores do tema
COR_FUNDO = RgbColor(26, 26, 46)  # #1a1a2e
COR_DESTAQUE = RgbColor(233, 69, 96)  # #e94560
COR_AZUL = RgbColor(0, 217, 255)  # #00d9ff
COR_TEXTO = RgbColor(255, 255, 255)
COR_CINZA = RgbColor(150, 150, 150)

def set_slide_background(slide, color):
    """Define cor de fundo do slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Adiciona slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COR_FUNDO)

    # Logo
    logo = slide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(1))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "i9ON"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COR_DESTAQUE
    p.alignment = PP_ALIGN.CENTER

    # Título
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.8), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(10), Inches(0.7))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = COR_AZUL
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items):
    """Adiciona slide com título e bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COR_FUNDO)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COR_DESTAQUE

    # Linha sob título
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COR_DESTAQUE
    line.line.fill.background()

    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        p.text = f"▸ {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    # Rodapé
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENCIAL | i9ON Tecnologia - Proposta de Investimento"
    p.font.size = Pt(10)
    p.font.color.rgb = COR_CINZA

    return slide

def add_stats_slide(prs, title, stats):
    """Adiciona slide com estatísticas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COR_FUNDO)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COR_DESTAQUE

    # Stats em grid
    col_width = 3
    start_x = 0.5
    start_y = 1.5

    for i, (number, label) in enumerate(stats):
        col = i % 3
        row = i // 3

        x = start_x + col * col_width
        y = start_y + row * 2.2

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(15, 52, 96)
        box.line.color.rgb = RgbColor(15, 52, 96)

        # Número
        num_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(2.8), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COR_AZUL
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(2.8), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = COR_CINZA
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows):
    """Adiciona slide com tabela"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COR_FUNDO)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COR_DESTAQUE

    # Tabela
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * num_rows)).table

    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COR_DESTAQUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COR_TEXTO

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RgbColor(20, 35, 60)
            else:
                cell.fill.fore_color.rgb = RgbColor(26, 26, 46)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COR_TEXTO

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Adiciona slide com duas colunas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COR_FUNDO)

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COR_DESTAQUE

    # Coluna esquerda - título
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_AZUL

    # Coluna esquerda - conteúdo
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"▸ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(8)

    # Coluna direita - título
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_AZUL

    # Coluna direita - conteúdo
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"▸ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(8)

    return slide

def main():
    # Criar apresentação
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Capa
    slide = add_title_slide(prs, "Proposta de Parceria Estratégica", "Oportunidade de Investimento")
    # Adicionar data
    date_box = slide.shapes.add_textbox(Inches(0), Inches(5.5), Inches(10), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dezembro 2024 | CONFIDENCIAL"
    p.font.size = Pt(18)
    p.font.color.rgb = COR_CINZA
    p.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Quem Somos
    add_two_column_slide(prs,
        "Quem Somos",
        "i9ON Tecnologia",
        [
            "Especialistas em ERP TOTVS Protheus",
            "10 anos de mercado consolidado",
            "Consultoria + Produtos próprios",
            "Clientes de grande porte",
            "Equipe técnica qualificada"
        ],
        "Nossos Produtos",
        [
            "i9 PDV Posto - Postos de combustível",
            "i9 SmartPDV - Vendas mobile e totens",
            "i9 SmartFeed - Gestão de dados",
            "i9 SmartCount - Controle de inventário",
            "Consultoria especializada Protheus"
        ]
    )

    # SLIDE 3: Nossos Números
    add_stats_slide(prs, "Nossos Números", [
        ("R$ 2,5M", "Faturamento Anual (2025)"),
        ("R$ 1,37M", "Lucro Líquido (2025)"),
        ("54,9%", "Margem Líquida"),
        ("10 anos", "No Mercado"),
        ("+37%", "Crescimento Lucro"),
        ("ZERO", "Dívidas Bancárias")
    ])

    # SLIDE 4: Evolução Financeira
    add_table_slide(prs, "Evolução Financeira",
        ["Indicador", "2024", "2025", "2026 (Proj.)", "Crescimento"],
        [
            ["Receita Bruta", "R$ 2.156.375", "R$ 2.493.128", "R$ 2.513.157", "+16,5%"],
            ["Lucro Bruto", "R$ 1.631.317", "R$ 1.830.093", "R$ 1.902.793", "+16,6%"],
            ["Lucro Líquido", "R$ 1.127.436", "R$ 1.369.228", "R$ 1.547.447", "+37,3%"],
            ["Margem Líquida", "52,3%", "54,9%", "61,6%", "Crescente"]
        ]
    )

    # SLIDE 5: Por Que Investir
    add_content_slide(prs, "Por Que Investir na i9ON?", [
        "Empresa altamente lucrativa - Margem de 55% (top 5% do setor)",
        "Mercado em crescimento - ERP Protheus em expansão",
        "Produtos próprios - Receita recorrente e propriedade intelectual",
        "10 anos de mercado - Empresa consolidada, baixo risco",
        "Zero endividamento - Saúde financeira excepcional",
        "Potencial de 2x - Dobrar faturamento em 24 meses com investimento"
    ])

    # SLIDE 6: Valuation
    add_table_slide(prs, "Valuation da Empresa",
        ["Método", "Múltiplo", "Valor da Empresa", "Observação"],
        [
            ["Múltiplo de Lucro", "4x", "R$ 5.476.912", "Conservador"],
            ["Múltiplo de Lucro", "5x", "R$ 6.846.140", "Moderado"],
            ["Múltiplo de Lucro", "5,5x", "R$ 6.500.000", "RECOMENDADO"],
            ["Múltiplo de Receita", "3x", "R$ 7.479.384", "SaaS/Produtos"]
        ]
    )

    # SLIDE 7: Modelos de Entrada
    add_content_slide(prs, "Modelos de Entrada", [
        "MODELO A: Investidor Passivo - Apenas capital, recebe dividendos",
        "MODELO B: Parceiro Estratégico - Capital + participação em decisões",
        "MODELO C: Cliente-Investidor - Estrutura especial para clientes",
        "MODELO D: Contrato - Adiantamento de serviços sem diluição",
        "MODELO E: Vesting - Participação gradual conforme metas",
        "MODELO F: Smart Money (RECOMENDADO) - Capital + equipe para profissionalizar"
    ])

    # SLIDE 8: Smart Money
    add_two_column_slide(prs,
        "Modelo Recomendado: Smart Money",
        "O Que o Investidor Traz",
        [
            "Aporte de capital (R$ 500k - R$ 800k)",
            "Diretor de Operações (mentoria)",
            "Equipe financeira (controles)",
            "RH compartilhado",
            "Consultor de processos"
        ],
        "O Que a i9ON Ganha",
        [
            "Profissionalização em 12 meses",
            "Processos estruturados",
            "Governança corporativa",
            "Sócios focados no core business",
            "Base para escalar 2x"
        ]
    )

    # SLIDE 9: Proposta Smart Money
    add_table_slide(prs, "Proposta: Smart Money (18%)",
        ["Componente", "Valor", "Participação"],
        [
            ["Aporte em capital", "R$ 600.000", "10%"],
            ["Equipe/consultoria (12 meses)", "R$ 350.000", "5%"],
            ["Garantia de contrato (5 anos)", "Receita garantida", "3%"],
            ["TOTAL", "~R$ 950.000 + contrato", "18%"]
        ]
    )

    # SLIDE 10: Cenários de Participação
    add_table_slide(prs, "Cenários de Participação",
        ["Cenário", "Aporte", "Participação", "Dividendos Anuais"],
        [
            ["Conservador", "R$ 600k - R$ 700k", "10%", "~R$ 137.000"],
            ["Moderado", "R$ 900k - R$ 1M", "15%", "~R$ 205.000"],
            ["Recomendado", "R$ 950k + equipe", "18%", "~R$ 246.000"],
            ["Agressivo", "R$ 1,2M - R$ 1,5M", "20%", "~R$ 274.000"]
        ]
    )

    # SLIDE 11: Nova Composição
    add_table_slide(prs, "Nova Composição Societária (18%)",
        ["Sócio", "Antes", "Depois"],
        [
            ["Lee Chardes (CEO)", "65%", "53,3%"],
            ["Rodrigo Felippe (CFO)", "30%", "24,6%"],
            ["Diego Rodrigues (CTO)", "5%", "4,1%"],
            ["Investidor", "0%", "18%"]
        ]
    )

    # SLIDE 12: Benefícios para Investidor
    add_content_slide(prs, "Benefícios para o Investidor", [
        "Dividendos anuais de ~R$ 250.000 (18% de participação)",
        "Margem líquida excepcional de 55%",
        "Crescimento projetado de +50% em 24 meses",
        "Payback estimado em 4-5 anos",
        "Participação em empresa de tecnologia lucrativa",
        "Garantia de fornecedor crítico (se for cliente)",
        "Assento no conselho + influência no roadmap"
    ])

    # SLIDE 13: Governança
    add_two_column_slide(prs,
        "Governança e Proteções",
        "Oferecemos ao Investidor",
        [
            "Assento no Conselho",
            "Relatórios Mensais",
            "Tag-Along (proteção em venda)",
            "Direito de Preferência",
            "Acesso total à informação"
        ],
        "Esperamos do Investidor",
        [
            "Lock-up de 5 anos",
            "Non-compete (não investir em concorrentes)",
            "Confidencialidade",
            "Não-solicitação de funcionários",
            "Contribuição ativa (mentoria)"
        ]
    )

    # SLIDE 14: Próximos Passos
    add_content_slide(prs, "Próximos Passos", [
        "1. ALINHAMENTO - Discussão de expectativas e modelo preferido",
        "2. DUE DILIGENCE - Auditoria financeira e operacional",
        "3. TERM SHEET - Carta de intenções com termos principais",
        "4. FORMALIZAÇÃO - Contratos, alteração societária e registro"
    ])

    # SLIDE 15: Encerramento
    slide = add_title_slide(prs, "Obrigado!", "Estamos prontos para crescer juntos")

    # Adicionar contatos
    contatos = slide.shapes.add_textbox(Inches(0), Inches(5), Inches(10), Inches(1.5))
    tf = contatos.text_frame
    p = tf.paragraphs[0]
    p.text = "Lee Chardes - CEO | Rodrigo Felippe - CFO | Diego Rodrigues - CTO"
    p.font.size = Pt(18)
    p.font.color.rgb = COR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "CONFIDENCIAL - Dezembro 2024"
    p.font.size = Pt(14)
    p.font.color.rgb = COR_DESTAQUE
    p.alignment = PP_ALIGN.CENTER

    # Salvar
    prs.save('apresentacao_investidor.pptx')
    print("Apresentação salva como 'apresentacao_investidor.pptx'")

if __name__ == "__main__":
    main()
